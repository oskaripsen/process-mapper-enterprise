"""
File Extraction Utilities

Extracts text and images from various file formats:
- PDF (PyMuPDF/fitz)
- DOCX (python-docx)
- PPTX (python-pptx)
- Images (PNG, JPEG)
"""

import io
import base64
from typing import List, Tuple
from fastapi import UploadFile, HTTPException
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image


async def extract_text(file: UploadFile) -> str:
    """
    Extract all readable text from a file.
    
    Args:
        file: FastAPI UploadFile object
        
    Returns:
        Extracted text as a string
        
    Raises:
        HTTPException: If file type is unsupported or extraction fails
    """
    try:
        # Read file content
        content = await file.read()
        filename = file.filename.lower() if file.filename else ""
        
        # Determine file type and extract
        if filename.endswith('.pdf'):
            return _extract_text_from_pdf(content)
        elif filename.endswith('.docx'):
            return _extract_text_from_docx(content)
        elif filename.endswith('.pptx'):
            return _extract_text_from_pptx(content)
        elif filename.endswith(('.png', '.jpg', '.jpeg')):
            # Images don't have extractable text (unless OCR is used)
            # We'll handle images separately via vision API
            return ""
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {filename}. Supported types: .pdf, .docx, .pptx, .png, .jpeg"
            )
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error extracting text from {file.filename}: {str(e)}"
        )


async def extract_images(file: UploadFile) -> List[str]:
    """
    Extract all images from a file as base64-encoded strings.
    
    Args:
        file: FastAPI UploadFile object
        
    Returns:
        List of base64-encoded image strings
        
    Raises:
        HTTPException: If file type is unsupported or extraction fails
    """
    try:
        # Read file content
        content = await file.read()
        filename = file.filename.lower() if file.filename else ""
        
        # Determine file type and extract
        if filename.endswith('.pdf'):
            return _extract_images_from_pdf(content)
        elif filename.endswith('.docx'):
            return _extract_images_from_docx(content)
        elif filename.endswith('.pptx'):
            return _extract_images_from_pptx(content)
        elif filename.endswith(('.png', '.jpg', '.jpeg')):
            return _extract_image_file(content)
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {filename}. Supported types: .pdf, .docx, .pptx, .png, .jpeg"
            )
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error extracting images from {file.filename}: {str(e)}"
        )


# ==================== PDF Extraction ====================

def _extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF using PyMuPDF."""
    doc = fitz.open(stream=content, filetype="pdf")
    text_parts = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        text_parts.append(page.get_text())
    
    doc.close()
    return "\n\n".join(text_parts)


def _extract_images_from_pdf(content: bytes) -> List[str]:
    """Extract images from PDF as base64 strings."""
    doc = fitz.open(stream=content, filetype="pdf")
    images = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images()
        
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            # Convert to base64
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            
            # Determine image format
            image_ext = base_image["ext"]
            mime_type = f"image/{image_ext}"
            
            # Format for OpenAI API
            images.append(f"data:{mime_type};base64,{base64_image}")
    
    doc.close()
    return images


# ==================== DOCX Extraction ====================

def _extract_text_from_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    doc = Document(io.BytesIO(content))
    text_parts = []
    
    # Extract from paragraphs
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
    
    # Extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text_parts.append(row_text)
    
    return "\n\n".join(text_parts)


def _extract_images_from_docx(content: bytes) -> List[str]:
    """Extract images from DOCX as base64 strings."""
    doc = Document(io.BytesIO(content))
    images = []
    
    # Extract inline images from relationships
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                
                # Convert to base64
                base64_image = base64.b64encode(image_data).decode('utf-8')
                
                # Determine MIME type
                content_type = rel.target_part.content_type
                images.append(f"data:{content_type};base64,{base64_image}")
            except Exception as e:
                print(f"Warning: Could not extract image from DOCX: {str(e)}")
                continue
    
    return images


# ==================== PPTX Extraction ====================

def _extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    prs = Presentation(io.BytesIO(content))
    text_parts = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = [f"--- Slide {slide_num} ---"]
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
            
            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_text.append(row_text)
        
        if len(slide_text) > 1:  # More than just the slide header
            text_parts.append("\n".join(slide_text))
    
    return "\n\n".join(text_parts)


def _extract_images_from_pptx(content: bytes) -> List[str]:
    """Extract images from PPTX as base64 strings."""
    prs = Presentation(io.BytesIO(content))
    images = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                try:
                    image = shape.image
                    image_data = image.blob
                    
                    # Convert to base64
                    base64_image = base64.b64encode(image_data).decode('utf-8')
                    
                    # Determine MIME type
                    content_type = image.content_type
                    images.append(f"data:{content_type};base64,{base64_image}")
                except Exception as e:
                    print(f"Warning: Could not extract image from PPTX: {str(e)}")
                    continue
    
    return images


# ==================== Image File Extraction ====================

def _extract_image_file(content: bytes) -> List[str]:
    """Convert image file to base64 string."""
    try:
        # Verify it's a valid image
        image = Image.open(io.BytesIO(content))
        
        # Convert to base64
        base64_image = base64.b64encode(content).decode('utf-8')
        
        # Determine MIME type
        format_to_mime = {
            'JPEG': 'image/jpeg',
            'PNG': 'image/png',
            'JPG': 'image/jpeg'
        }
        mime_type = format_to_mime.get(image.format, 'image/jpeg')
        
        return [f"data:{mime_type};base64,{base64_image}"]
        
    except Exception as e:
        raise Exception(f"Invalid image file: {str(e)}")

